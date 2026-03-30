#!/usr/bin/env python3
"""Generate the FYP viva presentation (18 slides, 20 minutes).

Matches the Presentation Guide requirements exactly.

Usage:
    python create_presentation.py
    # Output: AFLHR_Lite_Viva.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colour palette
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x16, 0x21, 0x3A)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
TEXT_WHITE = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
RED = RGBColor(0xF8, 0x71, 0x71)
ORANGE = RGBColor(0xFB, 0xBF, 0x24)


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, top=0.5, left=0.8, size=36, color=TEXT_WHITE, bold=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(11.5), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    return txBox


def add_text(slide, text, top, left=0.8, width=11.5, size=18, color=TEXT_MUTED, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_after = Pt(8)
    return txBox


def add_bullets(slide, items, top, left=0.8, width=11.5, size=18, color=TEXT_MUTED):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide)
add_text(slide, "INFORMATICS INSTITUTE OF TECHNOLOGY", 1.0, size=14, color=TEXT_MUTED)
add_text(slide, "In Collaboration with University of Westminster", 1.4, size=14, color=TEXT_MUTED)
add_title(slide, "An Adaptive, Confidence-Weighted Verification\nFramework for Mitigating LLM Hallucinations", top=2.2, size=32)
add_text(slide, "AFLHR Lite \u2014 Cw-CONLI Hallucination Detection", 3.8, size=20, color=ACCENT)
add_text(slide, "Shaun Yogeshwaran  \u2022  w1912919", 5.0, size=18, color=TEXT_WHITE)
add_text(slide, "Supervised by Mr. John Sriskandarajah", 5.5, size=16, color=TEXT_MUTED)
add_text(slide, "BSc (Hons) Computer Science  \u2022  March 2026", 6.2, size=14, color=TEXT_MUTED)
add_notes(slide, "Introduce yourself, state your project title and supervisor. ~30 seconds.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Agenda")
items = [
    "Problem Statement & Research Gap",
    "Research Aim & Questions",
    "System Design & Architecture",
    "Implementation & Technology Stack",
    "AI/ML Pipeline Details",
    "Live Demonstration",
    "Testing & Evaluation Results",
    "Critical Evaluation Outcomes",
    "Contributions, Novelty & Limitations",
    "Conclusion & Skills Acquired",
]
add_bullets(slide, items, 1.8, size=20, color=TEXT_WHITE)
add_notes(slide, "Just show the agenda, don't spell out each item. ~15 seconds.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: Problem Statement
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Problem Statement")
add_text(slide, "Large Language Models generate fluent but factually unsupported content", 1.8, size=22, color=TEXT_WHITE)
items = [
    "LLMs hallucinate \u2014 producing confident but incorrect information (Ji et al., 2023)",
    "Existing NLI verification uses uniform thresholds regardless of evidence quality",
    "Static thresholds over-flag well-supported responses and under-flag poorly-supported ones",
    "Summarisation tasks break when premises exceed 512 tokens (RoBERTa limit)",
    "No existing framework adapts verification stringency to retrieval confidence",
]
add_bullets(slide, items, 2.8, size=18)
add_notes(slide, "Explain the core problem: LLMs hallucinate, and existing verification doesn't adapt to evidence quality. Cite Ji et al. 2023 and Williams et al. 2018. ~1.5 minutes.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: Research Gap (MOST IMPORTANT)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Research Gap", color=ACCENT)
add_text(slide, "This is the core motivation for the project", 1.6, size=16, color=ORANGE)
items = [
    "Existing work (CoNLI, SGIC, Chain-of-Verification) uses fixed verification thresholds",
    "No framework adjusts NLI stringency based on retrieval confidence quality",
    "512-token truncation breaks NLI for long-context tasks (summarisation)",
    "Whole-response NLI misses partial hallucinations within multi-claim responses",
    "Gap: Need adaptive, confidence-weighted verification that handles variable evidence quality",
]
add_bullets(slide, items, 2.4, size=18)
add_notes(slide, "THIS IS THE MOST IMPORTANT SLIDE per the presentation guide. Clearly explain: no existing framework adapts NLI threshold based on retrieval confidence. Cite CoNLI, SGIC (Chen et al. 2025), Chain-of-Verification (Dhuliawala et al. 2024). ~2 minutes.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: Project Aim & Research Questions
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Research Aim & Questions")
add_text(slide, "Aim: To design, implement, and evaluate a confidence-weighted NLI verification framework that dynamically adjusts hallucination detection thresholds based on retrieval confidence.", 1.8, size=18, color=TEXT_WHITE)
items = [
    "RQ1: Does confidence-weighted thresholding improve hallucination detection F1?",
    "RQ2: Does Cw-CONLI reduce over-flagging of correct responses?",
    "RQ3: Which weighting function (tiered, sqrt, sigmoid) performs best?",
]
add_bullets(slide, items, 3.2, size=18, color=ACCENT)
add_notes(slide, "State the aim clearly, then present the three RQs. These drive the entire evaluation. ~1 minute.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: System Design
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "System Architecture")
add_text(slide, "[INSERT: Architecture diagram from thesis — Figure 3.1 or similar]", 2.0, size=16, color=ORANGE)
items = [
    "Layer 1: RAG Retrieval \u2014 FAISS vector search with BGE/MiniLM embeddings",
    "Layer 2: LLM Generation \u2014 Llama-3.1-8B via Groq API",
    "Layer 3: NLI Verification \u2014 RoBERTa-large-MNLI (whole/windowed/decomposed)",
    "Layer 4: Adaptive Verdict \u2014 Cw-CONLI threshold (tiered, sqrt, sigmoid)",
]
add_bullets(slide, items, 3.5, size=18)
add_notes(slide, "Walk through the 4-layer pipeline. Point to the architecture diagram. Explain how data flows from query to verdict. ~1.5 minutes.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: Implementation & Tech Stack
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Technology Stack")
items = [
    "Backend: Python 3.10, FastAPI, FAISS, HuggingFace Transformers",
    "Frontend: React 18, Vite 5, Framer Motion, Recharts",
    "ML Models: BGE-small-en-v1.5 (embeddings), RoBERTa-large-MNLI (NLI), Llama-3.1-8B (generation)",
    "Evaluation: HaluEval benchmark (20,000 samples), McNemar\u2019s test",
    "Infrastructure: CPU-only (M4 MacBook Pro), Makefile automation",
]
add_bullets(slide, items, 1.8, size=18, color=TEXT_WHITE)
add_text(slide, "[INSERT: Technology stack diagram]", 5.5, size=16, color=ORANGE)
add_notes(slide, "Justify key technology choices. Why FastAPI? Why FAISS? Why CPU-only? ~1 minute.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: AI/ML Model Details
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "AI/ML Pipeline Details")
items = [
    "Dataset: HaluEval (Li et al. 2023) \u2014 10K QA + 10K Summarisation samples",
    "Split: 70% dev / 30% test (seed=42), stratified by task",
    "Precompute: Retrieval + NLI scores cached once, then sweep thresholds (fast tuning)",
    "v2 Upgrades: Sliding-window NLI (400-token windows, 200-stride), claim decomposition, BGE embeddings",
    "Calibration: Temperature scaling investigated \u2014 T=10.0 at boundary (documented negative result)",
]
add_bullets(slide, items, 1.8, size=18, color=TEXT_WHITE)
add_text(slide, "Three-condition experiment: C1 (RAG-only) \u2192 C2 (static CONLI) \u2192 C3 (Cw-CONLI)", 5.2, size=20, color=ACCENT)
add_notes(slide, "Explain the experiment design. Precompute-then-sweep is a key engineering decision. Explain the v2 improvements. ~1.5 minutes.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: Demo
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Live Demonstration", size=40, color=GREEN)
add_text(slide, "Showing the AFLHR Lite React interface", 2.5, size=24, color=TEXT_WHITE)
items = [
    "Verify a factually supported query (Westminster founding)",
    "Detect a hallucination (off-topic query)",
    "v2 mode: per-claim decomposition breakdown",
    "Threshold tuning: pivot adjustment",
    "Batch exploration across domains",
]
add_bullets(slide, items, 3.5, size=20)
add_notes(slide, "LIVE DEMO. Have the system running (make start). Show: (1) verified query, (2) hallucination, (3) v2 per-claim, (4) threshold slider, (5) Explore page batch. ~3 minutes.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: Testing & Evaluation
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Testing & Evaluation Results")
items = [
    "Unit + Integration Testing: 36/36 passed (100%)",
    "API-Level Functional Testing: 18/18 passed (100%)",
    "Non-Functional Testing: 10/10 passed (100%) \u2014 avg response 405ms (v1), 690ms (v2)",
    "Model Evaluation: C2 F1 = 0.6988, C3 Tiered F1 = 0.6998 (combined test)",
    "QA: C3 Sqrt reduces over-flagging from 12.2% to 9.5%",
    "McNemar\u2019s p = 1.0 (standard) \u2014 C3 and C2 converge on this benchmark",
]
add_bullets(slide, items, 1.8, size=18, color=TEXT_WHITE)
add_notes(slide, "Present testing comprehensively: unit, integration, API-level, non-functional, and model evaluation. Highlight the null result honestly. ~1.5 minutes.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11: Critical Evaluation Outcomes
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Critical Evaluation Outcomes")
add_text(slide, "6 evaluators: 4 technical experts + 2 domain experts", 1.6, size=16, color=TEXT_MUTED)
items = [
    "Transparency rated highest: 4.5/5 \u2014 full exposure of scores and reasoning",
    "Deployment feasibility: 4.2/5 \u2014 CPU-only, sub-second, Makefile setup",
    "Detection accuracy: 3.8/5 \u2014 strong on QA, weaker on summarisation",
    "Algorithmic novelty: 3.2/5 \u2014 sound concept, limited impact on this benchmark",
    "Scalability: 3.0/5 \u2014 in-memory FAISS limits KB size",
    "Key feedback: \"Honest null result reporting is a strength, not a weakness\" (DE-1)",
]
add_bullets(slide, items, 2.2, size=18)
add_notes(slide, "Summarise expert evaluation. Be honest about ratings. The transparency and honest reporting were highlighted as strengths. ~1 minute.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12: Contribution to Body of Knowledge
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Contribution to Body of Knowledge")
add_text(slide, "Problem Domain:", 1.8, size=20, color=ACCENT, bold=True)
items_prob = [
    "Cw-CONLI framework: first system to weight NLI thresholds by retrieval confidence",
    "Realistic experiment shows significant benefit (p = 1.40\u00d710\u207b\u2075) when retrieval varies",
    "v2 engineering: windowed NLI fixes summarisation, decomposition catches partial hallucinations",
]
add_bullets(slide, items_prob, 2.4, size=17)
add_text(slide, "Research Domain:", 4.2, size=20, color=ACCENT, bold=True)
items_res = [
    "Reproducible three-condition experiment with dev/test split and McNemar\u2019s test",
    "Honest null result: C3 \u2248 C2 on standard benchmark (identifies boundary conditions)",
    "Documented negative result: calibration T=10.0 (NLI logits uncalibratable)",
]
add_bullets(slide, items_res, 4.8, size=17)
add_notes(slide, "Clearly state both contributions. Emphasise that the null result is informative — it tells the community WHEN adaptive thresholds help and when they don't. ~1.5 minutes.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13: Novelty
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Novelty")
items = [
    "Confidence-weighted NLI threshold: no prior work adapts verification stringency to retrieval quality",
    "Three weighting functions compared: tiered (step), sqrt (continuous), sigmoid (smooth)",
    "Sliding-window NLI: overlapping 400-token windows with stride 200, max aggregation",
    "Sentence-level claim decomposition: split \u2192 verify each \u2192 min score",
    "Dual-experiment design: standard per-sample + realistic shared-index retrieval",
]
add_bullets(slide, items, 1.8, size=18, color=TEXT_WHITE)
add_notes(slide, "Clearly articulate what is NEW. The confidence-weighted threshold is the core novelty. The v2 engineering (windowed NLI, decomposition) is the practical novelty. ~1 minute.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14: Limitations & Future Work
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Limitations & Future Work")
add_text(slide, "Limitations:", 1.8, size=20, color=RED, bold=True)
items_lim = [
    "Single benchmark (HaluEval) \u2014 tight retrieval clustering limits C3 differentiation",
    "Single NLI model (RoBERTa-MNLI) \u2014 not fine-tuned for hallucination detection",
    "Summarisation over-flagging remains ~98-99% despite windowed NLI",
    "CPU-only inference \u2014 practical but slower than GPU deployment",
]
add_bullets(slide, items_lim, 2.5, size=17)
add_text(slide, "Future Work:", 4.8, size=20, color=GREEN, bold=True)
items_fut = [
    "Evaluate on datasets with natural retrieval diversity (Natural Questions, TriviaQA)",
    "NLI ensembles and hallucination-specific fine-tuning",
    "Weighted window aggregation for improved summarisation handling",
    "Learn sigmoid parameters from data (k, pivot) instead of grid search",
]
add_bullets(slide, items_fut, 5.4, size=17)
add_notes(slide, "Be honest about limitations. Then show awareness of how to address them. ~1 minute.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15: Conclusion
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Conclusion")
items = [
    "Cw-CONLI provides a principled framework for adaptive hallucination detection",
    "On standard benchmarks: C3 \u2248 C2 (null result \u2014 retrieval scores too clustered)",
    "Under realistic retrieval: C3 significantly reduces over-flagging (100% \u2192 44.9%, p = 1.4\u00d710\u207b\u2075)",
    "Primary contribution is v2 engineering: windowed NLI, decomposition, BGE embeddings",
    "All four research objectives achieved; three RQs answered with nuance",
    "Framework is deployable on consumer hardware with sub-second response times",
]
add_bullets(slide, items, 1.8, size=18, color=TEXT_WHITE)
add_notes(slide, "Summarise the key takeaways. Lead with the honest finding, then the practical contribution. ~1 minute.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16: New Skills Acquired
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "New Skills Acquired")
items = [
    "Natural Language Inference (NLI) \u2014 RoBERTa-MNLI, entailment classification",
    "Vector similarity search \u2014 FAISS indexing, cosine similarity, embedding models",
    "LLM API integration \u2014 Groq API, prompt engineering, temperature control",
    "Retrieval-Augmented Generation (RAG) pipeline design",
    "Statistical hypothesis testing \u2014 McNemar\u2019s test for paired classifier comparison",
    "React frontend development \u2014 Framer Motion animations, Vite build system",
    "Experiment design \u2014 precompute-then-sweep, dev/test split methodology",
]
add_bullets(slide, items, 1.8, size=18, color=TEXT_WHITE)
add_notes(slide, "These are skills NOT from the curriculum. Emphasise NLI, FAISS, RAG, and statistical testing as particularly novel. ~45 seconds.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17: Use of Existing Skills
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Application of Existing Skills")
items = [
    "Python programming \u2014 modular OOP design (AFLHREngine class), FastAPI REST API",
    "Software engineering \u2014 Git version control, Makefile automation, CI patterns",
    "Machine learning fundamentals \u2014 train/dev/test splits, F1/precision/recall metrics",
    "Data analysis \u2014 pandas, matplotlib for experiment analysis and visualisation",
    "Research methodology \u2014 literature review, Saunders\u2019 research onion, mixed methods",
    "Technical writing \u2014 structured academic report with Harvard referencing",
]
add_bullets(slide, items, 1.8, size=18, color=TEXT_WHITE)
add_notes(slide, "These are skills FROM the degree programme applied to the project. ~45 seconds.")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18: Summary
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Summary", color=ACCENT)
items = [
    "Built AFLHR Lite: a complete RAG + NLI hallucination detection pipeline",
    "Proposed Cw-CONLI: confidence-weighted thresholds for adaptive verification",
    "Evaluated on 20,000 samples: honest null result on standard, significant on realistic",
    "v2 engineering fixes real problems: windowed NLI, claim decomposition, BGE embeddings",
    "CPU-deployable, sub-second, fully reproducible",
]
add_bullets(slide, items, 2.0, size=20, color=TEXT_WHITE)
add_text(slide, "Thank you. Questions?", 5.5, size=28, color=ACCENT, bold=True)
add_notes(slide, "End with a clear summary of contributions. Open for questions. Thank the panel.")

# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
output = "AFLHR_Lite_Viva.pptx"
prs.save(output)
print(f"Saved {output} ({len(prs.slides)} slides)")
print("\nNOTE: You need to add diagrams manually:")
print("  - Slide 6: System architecture diagram (from thesis Figure 3.1)")
print("  - Slide 7: Technology stack diagram")
print("  - Slide 9: Prepare live demo (make start)")
